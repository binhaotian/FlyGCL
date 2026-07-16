from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"
PPTX_PATH = OUT_DIR / "FlyPrompt_LibContinual_复现汇报.pptx"
SCRIPT_PATH = OUT_DIR / "录制讲稿_10分钟以内.md"

FIG_DIR = ROOT / "docs/libcontinual_reproduction/json_experiment_results/figures"
ANALYSIS_DIR = ROOT / "docs/libcontinual_reproduction/analysis"

FONT = "Microsoft YaHei"
TITLE = RGBColor(24, 40, 72)
TEXT = RGBColor(45, 50, 60)
MUTED = RGBColor(105, 116, 132)
ACCENT = RGBColor(24, 112, 160)
ACCENT_2 = RGBColor(198, 88, 48)
BG = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if hasattr(shape, "line"):
        shape.line.color.rgb = color


def set_run(run, size=20, color=TEXT, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, x, y, w, h, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.35, 12.1, 0.55, title, size=27, color=TITLE, bold=True)
    if subtitle:
        add_textbox(slide, 0.58, 0.92, 11.8, 0.35, subtitle, size=12.5, color=MUTED)
    line = slide.shapes.add_shape(1, Inches(0.58), Inches(1.28), Inches(12.1), Inches(0.02))
    set_fill(line, ACCENT)


def add_bullets(slide, x, y, w, h, items, size=17, color=TEXT, gap=0.92):
    for idx, item in enumerate(items):
        yy = y + idx * gap
        dot = slide.shapes.add_shape(9, Inches(x), Inches(yy + 0.08), Inches(0.12), Inches(0.12))
        set_fill(dot, ACCENT if idx % 2 == 0 else ACCENT_2)
        box = add_textbox(slide, x + 0.25, yy, w - 0.25, 0.65, item, size=size, color=color)
        box.text_frame.word_wrap = True


def add_card(slide, x, y, w, h, title, body, accent=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, WHITE)
    shape.line.color.rgb = RGBColor(216, 224, 232)
    stripe = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(stripe, accent)
    add_textbox(slide, x + 0.22, y + 0.15, w - 0.35, 0.33, title, size=15, color=TITLE, bold=True)
    add_textbox(slide, x + 0.22, y + 0.58, w - 0.35, h - 0.7, body, size=12.8, color=TEXT)


def add_image(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_table(slide, x, y, w, h, rows, cols, data, font_size=12):
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c in range(cols):
        table.columns[c].width = Inches(w / cols)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(232, 239, 246) if r == 0 else WHITE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, size=font_size, color=TITLE if r == 0 else TEXT, bold=(r == 0))
    return table


def add_footer(slide, num):
    add_textbox(slide, 10.85, 7.12, 1.75, 0.25, f"{num:02d} / 10", size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)


def build_ppt():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_textbox(slide, 0.8, 0.72, 11.6, 0.72, "FlyPrompt 论文复现与 LibContinual 代码移植", 30, TITLE, True)
    add_textbox(slide, 0.85, 1.55, 10.7, 0.42, "代码 + PPT + 10 分钟以内录制演示材料", 17, ACCENT, True)
    add_card(slide, 0.85, 2.35, 3.7, 1.15, "任务目标", "将 FlyPrompt 从 FlyGCL 迁移到 LibContinual，并完成可追溯复现实验报告。")
    add_card(slide, 4.85, 2.35, 3.7, 1.15, "核心方法", "Prompt experts + REAR 解析路由 + TE2 多时间尺度 EMA heads。", ACCENT_2)
    add_card(slide, 8.85, 2.35, 3.7, 1.15, "最终产物", "GitHub 代码、中文复现报告、JSON 证据、图表与演示路线。")
    add_image(slide, FIG_DIR / "paper_vs_lib_sup21k_auc_last.png", 0.95, 4.05, 5.7, 2.55)
    add_image(slide, FIG_DIR / "paper_vs_lib_component_ablation_auc.png", 6.85, 4.05, 5.55, 2.55)
    add_footer(slide, 1)

    # 2
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "复现问题：为什么 FlyPrompt 适合 GCL？", "General Continual Learning 场景下，任务边界不可靠、类别会重复出现、数据是单遍流式到达。")
    add_bullets(slide, 0.75, 1.75, 5.65, 4.5, [
        "传统 task/class incremental 假设任务边界较清晰，GCL 更接近真实在线数据流。",
        "Si-Blurry 同时包含 disjoint classes 与 blurry classes，要求模型处理类别重现。",
        "冻结预训练 ViT 可以降低训练成本，但需要可靠的 expert routing 和稳定分类头。",
        "FlyPrompt 的核心价值是用轻量组件提升在线持续学习稳定性。"
    ], size=16.3, gap=0.88)
    add_card(slide, 7.0, 1.75, 5.0, 1.0, "难点 1：路由", "没有可靠任务标签时，如何为输入样本选择合适 expert？")
    add_card(slide, 7.0, 3.05, 5.0, 1.0, "难点 2：遗忘", "每个 expert 只看到有限、不均衡样本时，如何减少短期偏置？", ACCENT_2)
    add_card(slide, 7.0, 4.35, 5.0, 1.0, "复现目标", "不是简单跑通代码，而是在 LibContinual 生命周期中重建算法语义。")
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "方法结构：Prompt Expert + REAR + TE2", "迁移实现保留论文中的三个关键机制。")
    add_card(slide, 0.75, 1.7, 3.65, 4.55, "Prompt Expert Bank", "冻结 ViT backbone，只训练每个 expert 的 prompt token 和分类头。\n\n优点：参数高效，避免全模型在线更新成本。")
    add_card(slide, 4.85, 1.7, 3.65, 4.55, "REAR Router", "使用随机扩展特征维护 G/Q 统计量，通过 ridge regression 解析求解 router。\n\n优点：路由轻量、稳定，不依赖任务标签。", ACCENT_2)
    add_card(slide, 8.95, 1.7, 3.65, 4.55, "TE2 Heads", "为每个 expert 维护 online head 和多个 EMA heads，推理时做 temporal ensemble。\n\n优点：兼顾快速适应和长期平滑。")
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "代码移植：按 LibContinual 接口重建", "迁移不是直接调用 FlyGCL，而是适配模型、数据、训练和评估边界。")
    add_table(slide, 0.72, 1.55, 11.85, 4.45, 5, 3, [
        ["模块", "迁移内容", "关键文件"],
        ["模型层", "PromptBank、RearRouter、FlyPrompt 组合", "models/flyprompt.py"],
        ["训练层", "online_iter 多步更新与 EMA 更新", "methods/flyprompt.py"],
        ["Trainer", "FlyPrompt 专用 observe_with_optimizer 分支", "methods/_trainer.py"],
        ["报告证据", "JSON 归档、聚合表格、图像对比", "docs/libcontinual_reproduction/"],
    ], font_size=11.3)
    add_textbox(slide, 0.85, 6.25, 11.2, 0.55, "代码演示时重点打开上述 4 个位置，能覆盖“做了什么、怎么跑、结果在哪里、如何追溯”。", 15, ACCENT, True)
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "实验协议与证据链", "主实验对齐 CIFAR-100 / Sup-21K / Si-Blurry 默认设置，五个随机种子。")
    add_bullets(slide, 0.75, 1.55, 5.9, 4.6, [
        "数据流：5 sessions，r_D=50%，r_B=10%，balanced split。",
        "模型：ViT-B/16 Sup-21K，backbone 冻结，prompt length=5。",
        "训练：Adam，lr=0.005，batch size=64，online_iter=3，AMP。",
        "指标：A_auc、A_last、A_avg、F_last、BWT_last。"
    ], size=15.8, gap=0.82)
    add_card(slide, 7.0, 1.55, 5.15, 1.05, "主报告", "docs/libcontinual_reproduction/FULL_REPRODUCTION_REPORT.md")
    add_card(slide, 7.0, 2.85, 5.15, 1.05, "结果证据", "json_experiment_results/local_reference_json/all_local_json/")
    add_card(slide, 7.0, 4.15, 5.15, 1.05, "表格映射", "json_experiment_results/TABLE_TO_JSON_MAP.md", ACCENT_2)
    add_card(slide, 7.0, 5.45, 5.15, 1.05, "图表输出", "json_experiment_results/figures/")
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "主结果：论文与 LibContinual 迁移版对比", "迁移版在三个数据集上均略低于论文，但差距稳定。")
    add_image(slide, FIG_DIR / "paper_vs_lib_sup21k_auc_last.png", 0.7, 1.55, 6.15, 3.1)
    add_image(slide, FIG_DIR / "paper_vs_lib_sup21k_gap.png", 6.95, 1.55, 5.65, 3.1)
    add_table(slide, 0.9, 5.02, 11.6, 1.35, 4, 5, [
        ["数据集", "论文 A_auc", "迁移 A_auc", "论文 A_last", "迁移 A_last"],
        ["CIFAR-100", "83.24", "81.47 ± 0.61", "86.76", "84.59 ± 0.37"],
        ["ImageNet-R", "56.58", "55.12 ± 0.74", "55.27", "53.61 ± 0.68"],
        ["CUB-200", "70.64", "69.08 ± 0.92", "73.40", "71.46 ± 0.81"],
    ], font_size=10.2)
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "消融结果：关键组件趋势被复现", "完整 FlyPrompt 最好，去掉 prompt / REAR / EMA 均会下降。")
    add_image(slide, FIG_DIR / "paper_vs_lib_component_ablation_auc.png", 0.85, 1.5, 6.0, 3.75)
    add_table(slide, 7.05, 1.62, 5.35, 3.2, 5, 3, [
        ["设置", "迁移 A_auc", "迁移 A_last"],
        ["No prompt, no EMA", "69.82", "71.15"],
        ["No REAR, EMA", "80.28", "82.20"],
        ["REAR + prompt, no EMA", "80.15", "82.71"],
        ["Full FlyPrompt", "81.47", "84.59"],
    ], font_size=10.6)
    add_textbox(slide, 1.0, 5.75, 11.6, 0.55, "结论：迁移版不仅接近主结果，更重要的是保留了论文中“组件组合贡献”的相对关系。", 15, ACCENT, True)
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "在线曲线与异常处理", "GCL 评估会随 seen-class 范围扩大而阶段性变难，因此曲线下跳不是简单异常。")
    add_image(slide, ANALYSIS_DIR / "online_gcl_curve_seed1.png", 0.75, 1.55, 6.15, 3.95)
    add_bullets(slide, 7.2, 1.65, 5.0, 3.8, [
        "session 边界后 accuracy 短期回落，后续逐步恢复。",
        "五个 seed 全部纳入统计，没有手动剔除低 seed。",
        "AMP 和 transform 差异作为 ablation 记录，不混入主结果。",
        "所有主要数值可从 JSON、聚合表和 SHA256 追溯。"
    ], size=14.5, gap=0.78)
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "录制中的代码演示路线", "建议控制在 3 分钟以内，只展示核心路径，不现场跑完整实验。")
    add_bullets(slide, 0.75, 1.55, 11.4, 4.8, [
        "1. README：展示仓库入口和复现报告链接。",
        "2. models/flyprompt.py：展示 PromptBank、RearRouter、EMA heads。",
        "3. methods/flyprompt.py 与 methods/_trainer.py：展示 online_iter 多步更新和 trainer 接入。",
        "4. docs/libcontinual_reproduction/FULL_REPRODUCTION_REPORT.md：展示主报告。",
        "5. json_experiment_results/：展示 JSON 证据、聚合表格和图表。"
    ], size=16, gap=0.78)
    add_card(slide, 0.9, 6.05, 11.25, 0.65, "不建议现场执行完整训练", "10 分钟视频里只展示命令位置、配置文件和已归档结果；完整训练耗时较长。", ACCENT_2)
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, BG)
    add_title(slide, "结论与提交状态", "当前版本适合作为“代码 + PPT + 演示视频”的提交材料。")
    add_card(slide, 0.85, 1.55, 3.75, 2.05, "完成内容", "FlyPrompt 核心算法已迁移；报告、图表、JSON 证据已归档；GitHub 已更新。")
    add_card(slide, 4.85, 1.55, 3.75, 2.05, "复现结论", "主结果差距约 1-2 个百分点；组件趋势、EMA/mask/router 结论与论文一致。", ACCENT_2)
    add_card(slide, 8.85, 1.55, 3.75, 2.05, "局限说明", "部分论文扩展行无本地 JSON 映射；成本表仍依赖硬件和测量口径。")
    add_bullets(slide, 1.05, 4.35, 10.6, 1.8, [
        "提交仓库：https://github.com/binhaotian/FlyGCL",
        "最新提交：见 GitHub main 分支提交记录",
        "建议视频结构：PPT 约 7 分钟 + 代码演示约 3 分钟。"
    ], size=15.5, gap=0.58)
    add_footer(slide, 10)

    prs.save(PPTX_PATH)


def build_script():
    content = """# FlyPrompt 复现汇报录制讲稿（10 分钟以内）

## 时间分配

| 部分 | 建议时长 | 内容 |
| --- | ---: | --- |
| PPT 第 1-3 页 | 2 分钟 | 介绍任务、GCL 背景和 FlyPrompt 方法结构 |
| PPT 第 4-5 页 | 2 分钟 | 说明代码移植方式、实验协议和证据链 |
| PPT 第 6-8 页 | 3 分钟 | 展示主结果、消融趋势和在线曲线解释 |
| PPT 第 9 页 | 2 分钟 | 切到代码，演示 README、核心实现和结果目录 |
| PPT 第 10 页 | 1 分钟 | 总结结论、提交状态和局限 |

## 逐页讲稿

### 1. 标题页
本次工作完成了 FlyPrompt 论文在 LibContinual 框架中的代码移植和复现整理。交付内容包括 GitHub 代码、中文复现报告、JSON 证据、结果图表，以及这个 10 分钟以内的 PPT 和代码演示。

### 2. 复现问题
FlyPrompt 面向 General Continual Learning。这个设置比传统持续学习更难，因为任务边界不可靠，类别会在数据流中重复出现，模型只能单遍在线学习。因此复现时不能只看最终准确率，还要关注数据流顺序、在线评估和遗忘指标。

### 3. 方法结构
FlyPrompt 的核心是三部分：Prompt Expert Bank 负责参数高效适配，REAR Router 用随机扩展和 ridge regression 做解析式 expert routing，TE2 Heads 用多个 EMA 分类头平滑在线训练波动。迁移版保留了这三个机制。

### 4. 代码移植
这次不是直接调用 FlyGCL 源码，而是按 LibContinual 的接口重新实现。模型层在 `models/flyprompt.py`，训练逻辑在 `methods/flyprompt.py`，trainer 接入在 `methods/_trainer.py`，报告和证据在 `docs/libcontinual_reproduction/`。

### 5. 实验协议
主实验是 CIFAR-100 / Sup-21K / Si-Blurry，5 个 session，`r_D=50%`，`r_B=10%`，五个随机种子。指标包括 `A_auc`、`A_last`、`A_avg`、`F_last` 和 `BWT_last`。所有主要结果都有 JSON、聚合表和图表对应。

### 6. 主结果
从主结果看，LibContinual 迁移版在 CIFAR-100、ImageNet-R 和 CUB-200 上均略低于论文，大约 1 到 2 个百分点，但差距稳定，没有某个数据集完全失效。CIFAR-100/Sup-21K 上迁移版为 `A_auc=81.47±0.61`，`A_last=84.59±0.37`。

### 7. 消融结果
消融实验是判断迁移是否可信的关键。完整 FlyPrompt 最好；去掉 prompt 后下降明显；去掉 REAR 或 EMA 后处于中间。这说明迁移版不仅跑出了接近数值，也保留了论文中关键组件的相对贡献。

### 8. 在线曲线
在线曲线中有阶段性下跳，主要发生在新 session 开始附近。这不是简单异常，而是 seen-class 范围扩大后测试更难导致的。报告没有删点或平滑，五个 seed 都进入统计。

### 9. 代码演示
切到代码后，建议按这个顺序展示：

1. 打开 `README.md`，展示复现报告入口。
2. 打开 `models/flyprompt.py`，展示 PromptBank、RearRouter 和 EMA heads。
3. 打开 `methods/flyprompt.py`，展示 `online_iter` 多步更新。
4. 打开 `methods/_trainer.py`，展示 FlyPrompt 专用 trainer 分支。
5. 打开 `docs/libcontinual_reproduction/json_experiment_results/`，展示 JSON、聚合表和 figures。

不要现场跑完整训练，可以说明完整实验耗时较长，视频中只展示命令和结果证据。

### 10. 结论
最终结论是：FlyPrompt 的核心结构已经成功迁移到 LibContinual；主结果接近论文，关键消融趋势一致；当前仓库已经具备代码、中文报告、图表和可追溯 JSON 证据。局限是部分论文扩展行没有本地 JSON 映射，成本表也受硬件测量口径影响。
"""
    SCRIPT_PATH.write_text(content, encoding="utf-8")


if __name__ == "__main__":
    build_ppt()
    build_script()
    print(PPTX_PATH)
    print(SCRIPT_PATH)
